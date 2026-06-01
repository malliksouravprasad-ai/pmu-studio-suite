"""PowerPoint report generator for APP-008."""
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .deliverable_model import ReportConfig, SectionDef

_NAVY   = RGBColor(0x1F, 0x38, 0x64)
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT  = RGBColor(0xD9, 0xE1, 0xF2)
_SILVER = RGBColor(0xF2, 0xF2, 0xF2)

_W  = Inches(13.33)   # slide width
_H  = Inches(7.5)     # slide height
_BH = Inches(1.15)    # banner height


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def _blank_slide(prs: Presentation):
    """Add a blank slide with white background."""
    layout = prs.slide_layouts[6]   # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _WHITE
    return slide


def _add_banner(slide, title_text: str, subtitle_text: str = ""):
    """Add a navy banner at the top with white title text."""
    banner = slide.shapes.add_shape(1, 0, 0, _W, _BH)   # 1 = rectangle
    banner.fill.solid()
    banner.fill.fore_color.rgb = _NAVY
    banner.line.fill.background()

    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Remove existing runs then build
    tf.text = title_text
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = _WHITE

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size=11, bold=False, color=None, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.text = text
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def _add_ppt_table(slide, df: pd.DataFrame, top=None):
    """Add a formatted table to slide (max 12 data rows for readability)."""
    if df.empty:
        _add_text_box(slide, "(No data available.)",
                      Inches(0.5), _BH + Inches(0.2),
                      _W - Inches(1.0), Inches(0.5))
        return

    # Truncate for slide readability
    MAX_ROWS = 12
    is_truncated = len(df) > MAX_ROWS
    display_df = df.head(MAX_ROWS)

    n_rows = len(display_df) + 1   # +1 for header
    n_cols = len(display_df.columns)

    top = top or (_BH + Inches(0.2))
    left  = Inches(0.4)
    width = _W - Inches(0.8)
    height = _H - top - Inches(0.3)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Column widths (equal)
    col_w = width // n_cols
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    # Header
    for ci, col_name in enumerate(display_df.columns):
        cell = tbl.cell(0, ci)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY
        if cell.text_frame.paragraphs[0].runs:
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.color.rgb = _WHITE
            run.font.bold = True
            run.font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, (_, row) in enumerate(display_df.iterrows(), start=1):
        is_total = str(row.iloc[0]) == "TOTAL" if len(row) > 0 else False
        for ci, (col, val) in enumerate(row.items()):
            cell = tbl.cell(ri, ci)
            cell.text = str(val) if pd.notna(val) else ""
            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _LIGHT
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _SILVER
            if cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(8)
                if is_total:
                    run.font.bold = True

    if is_truncated:
        _add_text_box(
            slide,
            f"* Showing first {MAX_ROWS} rows. Full data in the Excel/Word report.",
            Inches(0.4), _H - Inches(0.35), _W - Inches(0.8), Inches(0.3),
            font_size=8, color=RGBColor(0x59, 0x59, 0x59))


def save_ppt_report(
    section_data: dict, config: ReportConfig, artifact_id: str
) -> tuple:
    prs = _new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    title_slide = _blank_slide(prs)
    # Full navy background
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _NAVY

    _add_text_box(title_slide, config.report_title,
                  Inches(0.8), Inches(2.0), _W - Inches(1.6), Inches(1.8),
                  font_size=36, bold=True, color=_WHITE)
    _add_text_box(title_slide,
                  "\n".join(filter(None, [
                      config.programme,
                      config.organization,
                      config.report_date,
                  ])),
                  Inches(0.8), Inches(3.9), _W - Inches(1.6), Inches(1.5),
                  font_size=14, color=RGBColor(0xBD, 0xD7, 0xEE))

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    agenda_slide = _blank_slide(prs)
    _add_banner(agenda_slide, "Agenda")
    sections_list = "\n".join(
        f"  {i}. {s.title}"
        for i, s in enumerate(config.sections, 1) if s.enabled)
    _add_text_box(agenda_slide, sections_list,
                  Inches(0.5), _BH + Inches(0.3), _W - Inches(1.0), _H - _BH - Inches(0.5),
                  font_size=14, color=_NAVY)

    # ── Section slides ────────────────────────────────────────────────────────
    for sdef in config.sections:
        if not sdef.enabled:
            continue
        slide = _blank_slide(prs)
        _add_banner(slide, sdef.title)

        if sdef.section_type == "narrative":
            text = sdef.narrative or "(No content provided.)"
            _add_text_box(slide, text,
                          Inches(0.5), _BH + Inches(0.3),
                          _W - Inches(1.0), _H - _BH - Inches(0.5),
                          font_size=13, color=_NAVY)

        else:
            df = section_data.get(sdef.section_id, pd.DataFrame())
            if sdef.narrative:
                _add_text_box(slide, sdef.narrative,
                              Inches(0.5), _BH + Inches(0.15),
                              _W - Inches(1.0), Inches(0.6),
                              font_size=10, color=RGBColor(0x40, 0x40, 0x40))
                top = _BH + Inches(0.8)
            else:
                top = _BH + Inches(0.2)
            _add_ppt_table(slide, df, top=top)

    # ── Last slide: Thank You ────────────────────────────────────────────────
    ty_slide = _blank_slide(prs)
    bg = ty_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _NAVY
    _add_text_box(ty_slide, "Thank You",
                  Inches(1.0), Inches(2.8), _W - Inches(2.0), Inches(1.5),
                  font_size=40, bold=True, color=_WHITE)
    _add_text_box(ty_slide,
                  f"{config.organization}  |  {config.report_date}",
                  Inches(1.0), Inches(4.4), _W - Inches(2.0), Inches(0.6),
                  font_size=14, color=RGBColor(0xBD, 0xD7, 0xEE))

    fname = f"{artifact_id}_report.pptx"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return fname, buf
